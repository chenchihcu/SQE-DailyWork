from __future__ import annotations

import argparse
import json
import os
import shutil
import sqlite3
import sys
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path
from typing import Any


REPO_ROOT = Path(__file__).resolve().parents[1]
SRC_ROOT = REPO_ROOT / "src"
for candidate in (SRC_ROOT, REPO_ROOT):
    text = str(candidate)
    if text not in sys.path:
        sys.path.insert(0, text)
VENV_SITE_PACKAGES = REPO_ROOT / ".venv" / "Lib" / "site-packages"
if VENV_SITE_PACKAGES.exists() and str(VENV_SITE_PACKAGES) not in sys.path:
    sys.path.append(str(VENV_SITE_PACKAGES))


@dataclass
class EvidenceRow:
    surface: str
    source_table: str
    source_field: str
    service_function: str
    output_artifact: str
    output_location: str
    expected: Any
    actual: Any
    status: str
    evidence_path: str


class EvidenceMatrix:
    def __init__(self) -> None:
        self.rows: list[EvidenceRow] = []

    def add(
        self,
        *,
        surface: str,
        source_table: str,
        source_field: str,
        service_function: str,
        output_artifact: str,
        output_location: str,
        expected: Any,
        actual: Any,
        evidence_path: Path | str,
        status: str | None = None,
    ) -> None:
        resolved_status = status or ("pass" if expected == actual else "not_pass")
        self.rows.append(
            EvidenceRow(
                surface=surface,
                source_table=source_table,
                source_field=source_field,
                service_function=service_function,
                output_artifact=output_artifact,
                output_location=output_location,
                expected=expected,
                actual=actual,
                status=resolved_status,
                evidence_path=str(evidence_path),
            )
        )

    def add_contains(
        self,
        *,
        surface: str,
        source_table: str,
        source_field: str,
        service_function: str,
        output_artifact: str,
        output_location: str,
        expected: str,
        actual: str,
        evidence_path: Path | str,
    ) -> None:
        self.add(
            surface=surface,
            source_table=source_table,
            source_field=source_field,
            service_function=service_function,
            output_artifact=output_artifact,
            output_location=output_location,
            expected=f"contains:{expected}",
            actual=actual,
            status="pass" if expected in actual else "not_pass",
            evidence_path=evidence_path,
        )

    def write(self, output_dir: Path) -> tuple[Path, Path]:
        json_path = output_dir / "evidence_matrix.json"
        md_path = output_dir / "evidence_matrix.md"
        data = [asdict(row) for row in self.rows]
        json_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

        lines = [
            "# SQE DailyWork Evidence Matrix",
            "",
            "| Status | Surface | Source | Service | Output | Expected | Actual | Evidence |",
            "| --- | --- | --- | --- | --- | --- | --- | --- |",
        ]
        for row in self.rows:
            source = f"{row.source_table}.{row.source_field}"
            output = f"{row.output_artifact}:{row.output_location}"
            lines.append(
                "| {status} | {surface} | `{source}` | `{service}` | `{output}` | {expected} | {actual} | `{evidence}` |".format(
                    status=row.status,
                    surface=_md(row.surface),
                    source=source,
                    service=row.service_function,
                    output=output,
                    expected=_md(row.expected),
                    actual=_md(row.actual),
                    evidence=row.evidence_path,
                )
            )
        md_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
        return json_path, md_path

    def exit_code(self) -> int:
        return 0 if all(row.status == "pass" for row in self.rows) else 1


def _md(value: Any) -> str:
    return str(value).replace("|", "\\|").replace("\n", " ")


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Run disposable DB data-to-report parity checks for SQE DailyWork."
    )
    parser.add_argument("--db-path", type=Path, default=REPO_ROOT / "data" / "sqe_v2.db")
    parser.add_argument("--output-dir", type=Path)
    parser.add_argument("--start-date", default="2026-06-01")
    parser.add_argument("--end-date", default="2026-06-30")
    parser.add_argument("--profile", choices=("focused", "full"), default="full")
    return parser.parse_args()


def _prepare_output_dir(raw: Path | None) -> Path:
    if raw is None:
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        raw = REPO_ROOT / "Outputs" / "audit" / stamp / "report-parity"
    output_dir = raw.expanduser().resolve()
    allowed = (REPO_ROOT / "Outputs" / "audit").resolve()
    if not str(output_dir).lower().startswith(str(allowed).lower()):
        raise ValueError(f"output-dir must stay under {allowed}")
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir


def _prepare_disposable_db(source: Path, output_dir: Path) -> Path:
    target = output_dir / "sqe_v2_audit.db"
    if source.exists():
        from database.backup import backup_sqlite_database

        backup_sqlite_database(source, target)
    else:
        target.parent.mkdir(parents=True, exist_ok=True)
        sqlite3.connect(target).close()
    return target


def _setup_environment(db_path: Path, output_dir: Path) -> None:
    os.environ["SQE_DB_PATH"] = str(db_path)
    os.environ["SQE_REQUIRE_DISPOSABLE_DB"] = "1"
    os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
    existing_pythonpath = [
        item for item in os.environ.get("PYTHONPATH", "").split(os.pathsep) if item
    ]
    os.environ["PYTHONPATH"] = os.pathsep.join(
        [str(SRC_ROOT), str(REPO_ROOT), str(VENV_SITE_PACKAGES), *existing_pythonpath]
    )
    (output_dir / "artifacts").mkdir(parents=True, exist_ok=True)


def _load_pdf_text(path: Path) -> tuple[str, str]:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        return "", f"not_verified: pypdf unavailable: {exc}"
    try:
        reader = PdfReader(str(path))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return text, "pass"
    except Exception as exc:
        return "", f"not_verified: PDF text extraction failed: {exc}"


def _pptx_text(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.extend(cell.text for cell in row.cells)
    return "\n".join(texts)


def _ensure_qapplication() -> object | None:
    try:
        from PySide6.QtWidgets import QApplication
    except Exception:
        return None
    return QApplication.instance() or QApplication([])


def _row_by_value(rows: list[dict], key: str, value: str) -> dict:
    for row in rows:
        if str(row.get(key) or "") == value:
            return row
    raise AssertionError(f"Could not find {key}={value}")


def _seed_and_audit(args: argparse.Namespace, output_dir: Path, matrix: EvidenceMatrix) -> dict[str, Any]:
    from openpyxl import Workbook, load_workbook

    from database import connection, repository
    from ncr.services import defect_service
    from ncr.services import export_service as ncr_export_service
    from ncr.services import stats_service as ncr_stats_service
    from services import attachment_manager, event_pdf_exporter, event_service
    from services import master_import_service
    from services.event import _query_service

    connection.initialize_database()
    _ensure_qapplication()
    stamp = datetime.now().strftime("%H%M%S")
    supplier_name = f"AUDIT供應商{stamp}"
    product_code = f"AUDIT-{stamp}"
    product_name = f"AUDIT產品{stamp}"
    imported_supplier = f"AUDIT匯入供應商{stamp}"
    imported_product_code = f"IMP-{stamp}"
    imported_product_name = f"AUDIT匯入產品{stamp}"
    artifacts_dir = output_dir / "artifacts"

    supplier_id = event_service.create_supplier(
        {
            "supplier_name": supplier_name,
            "contact_name": "王小明",
            "department": "品質",
            "phone": "02-12345678",
            "contact_email": f"audit-{stamp}@example.com",
        }
    )
    product_id = event_service.create_product(
        {
            "product_code": product_code,
            "product_name": product_name,
            "product_stage": "試產",
            "supplier_id": supplier_id,
        }
    )
    event_service.add_supplier_contact(
        supplier_id,
        {
            "contact_name": "李小華",
            "department": "工程",
            "phone": "02-87654321",
            "email": f"audit2-{stamp}@example.com",
            "is_primary": True,
        },
    )
    event_service.update_product(
        product_id,
        {
            "product_code": product_code,
            "product_name": product_name,
            "product_stage": "量產",
            "supplier_id": supplier_id,
            "stage_change_reason": "稽核驗證階段同步",
        },
    )

    suppliers = event_service.list_suppliers()
    products = event_service.list_products()
    contacts = event_service.list_supplier_contacts(supplier_id)
    stage_logs = event_service.list_product_stage_change_logs(product_id=product_id)
    matrix.add(
        surface="Master data supplier",
        source_table="suppliers",
        source_field="supplier_name",
        service_function="event_service.list_suppliers",
        output_artifact="service payload",
        output_location="suppliers[].supplier_name",
        expected=supplier_name,
        actual=_row_by_value(suppliers, "supplier_name", supplier_name)["supplier_name"],
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Master data contact primary",
        source_table="supplier_contacts",
        source_field="is_primary",
        service_function="event_service.list_supplier_contacts",
        output_artifact="service payload",
        output_location="contacts[].is_primary",
        expected=True,
        actual=any(c.get("contact_name") == "李小華" and bool(c.get("is_primary")) for c in contacts),
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Master data product stage",
        source_table="products",
        source_field="product_stage",
        service_function="event_service.list_products",
        output_artifact="service payload",
        output_location="products[].product_stage",
        expected="量產",
        actual=_row_by_value(products, "product_code", product_code)["product_stage"],
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Master data stage audit log",
        source_table="product_stage_change_logs",
        source_field="reason",
        service_function="event_service.list_product_stage_change_logs",
        output_artifact="service payload",
        output_location="stage_logs[0].reason",
        expected=True,
        actual=any("稽核驗證" in str(row.get("reason") or "") for row in stage_logs),
        evidence_path=output_dir / "evidence_matrix.json",
    )

    import_xlsx = artifacts_dir / "master_import_source.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["料號", "品名", "供應商", "產品階段"])
    ws.append([imported_product_code, imported_product_name, imported_supplier, "量產"])
    wb.save(import_xlsx)
    with connection.get_connection() as conn:
        preview = master_import_service.preview_product_master_import(conn, import_xlsx)
        result = master_import_service.apply_product_master_import(
            conn, preview, source_file=import_xlsx
        )
    matrix.add(
        surface="Master import preview",
        source_table="products",
        source_field="product_code",
        service_function="master_import_service.preview_product_master_import",
        output_artifact="service payload",
        output_location="preview.add_count",
        expected=1,
        actual=preview.add_count,
        evidence_path=import_xlsx,
    )
    matrix.add(
        surface="Master import apply",
        source_table="import_batches",
        source_field="status",
        service_function="master_import_service.apply_product_master_import",
        output_artifact="service payload",
        output_location="result.added_count",
        expected=1,
        actual=result.added_count,
        evidence_path=import_xlsx,
    )

    standalone_no = event_service.create_anomaly(
        {
            "anomaly_date": "2026-06-10",
            "supplier_id": supplier_id,
            "product_id": product_id,
            "problem_desc": f"{stamp} standalone anomaly input",
            "category": "原物料異常",
            "product_lot_no": f"LOT-{stamp}",
            "outsource_work_order": f"WO-{stamp}",
            "batch_qty": 25,
            "pending_items": "等待供應商回覆",
            "responsible_person": "SQE",
            "due_date": "2026-06-20",
            "quality_report_required": True,
        }
    )
    with connection.get_connection() as conn:
        standalone_row = conn.execute(
            "SELECT id FROM anomalies WHERE anomaly_no = ?", (str(standalone_no),)
        ).fetchone()
    standalone_id = str(standalone_row["id"])
    event_service.close_anomaly(
        standalone_id,
        "稽核驗證改善完成",
        closed_by="SQE",
        root_cause_category="物料/來料品質異常",
        closed_at="2026-06-18",
    )

    visit_id = event_service.create_visit(
        {
            "visit_date": "2026-06-11",
            "supplier_id": supplier_id,
            "product_id": product_id,
            "visitor_name": "SQE",
            "summary": f"{stamp} visit summary input",
            "work_order_no": f"VWO-{stamp}",
            "production_qty": 30,
            "defect_notes": [
                {
                    "defect_desc": f"{stamp} visit defect note input",
                    "improvement_desc": "",
                    "note": "待供應商補改善",
                }
            ],
            "tech_transfer": True,
            "tech_transfer_states": {
                "tech_transfer_doc": "yes",
                "carrier_requirement": "na",
                "dispensing_process": "yes",
                "functional_test": "no",
                "packaging_requirement": "yes",
            },
        }
    )
    pending_note = next(
        row
        for row in event_service.list_pending_visit_defect_notes(limit=50)
        if str(row.get("defect_desc") or "") == f"{stamp} visit defect note input"
    )
    confirmed = event_service.confirm_visit_defect_note_as_anomaly(
        str(pending_note["id"]),
        {"product_id": product_id, "responsible_person": "SQE", "due_date": "2026-06-25"},
    )
    linked = event_service.create_anomaly_with_visit_link(
        {
            "visit_id": visit_id,
            "sync_visit": False,
            "anomaly_date": "2026-06-12",
            "supplier_id": supplier_id,
            "product_id": product_id,
            "problem_desc": f"{stamp} linked anomaly input",
            "category": "製程異常",
            "batch_qty": 12,
            "responsible_person": "SQE",
            "due_date": "2026-06-26",
            "quality_report_required": False,
        }
    )

    image_src = artifacts_dir / "attachment_source.png"
    from PIL import Image as PilImage

    PilImage.new("RGB", (32, 24), "#2F80ED").save(image_src)
    attachment_manager.import_single_anomaly_attachment(
        standalone_id, image_src, "audit_attachment.png"
    )
    attachment_manager.set_anomaly_captions(
        standalone_id, {"audit_attachment.png": "稽核附件說明"}
    )

    events = event_service.list_events({"supplier": supplier_name, "yyyymm": "202606"})
    range_events = event_service.list_events_by_range(args.start_date, args.end_date)
    audit_range_events = [
        row for row in range_events if str(row.get("supplier_name") or "") == supplier_name
    ]
    totals, ranking = _query_service.summarize_range_events(audit_range_events)
    matrix.add(
        surface="Supplier events list payload",
        source_table="anomalies",
        source_field="anomaly_no",
        service_function="event_service.list_events",
        output_artifact="service payload",
        output_location="events[].ref_no",
        expected=str(standalone_no),
        actual=_row_by_value(events, "ref_no", str(standalone_no))["ref_no"],
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Supplier events range summary",
        source_table="anomalies+visits",
        source_field="event_date",
        service_function="event_service.list_events_by_range/summarize_range_events",
        output_artifact="service payload",
        output_location="totals.total_anomalies",
        expected=3,
        actual=totals["total_anomalies"],
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Supplier visit defect conversion",
        source_table="visit_defect_notes",
        source_field="confirmed_anomaly_id",
        service_function="event_service.confirm_visit_defect_note_as_anomaly",
        output_artifact="service payload",
        output_location="confirmed.anomaly_id",
        expected=True,
        actual=bool(confirmed.get("anomaly_id")),
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Supplier linked anomaly",
        source_table="anomalies",
        source_field="visit_id",
        service_function="event_service.create_anomaly_with_visit_link",
        output_artifact="service payload",
        output_location="linked.visit_action",
        expected="linked",
        actual=linked.get("visit_action"),
        evidence_path=output_dir / "evidence_matrix.json",
    )

    monthly_xlsx = artifacts_dir / "event_monthly.xlsx"
    ok, msg = event_service.export_monthly_excel(str(monthly_xlsx), "202606")
    monthly_wb = load_workbook(monthly_xlsx, data_only=False)
    monthly_text = "\n".join(
        str(cell.value or "")
        for sheet in monthly_wb.worksheets
        for row in sheet.iter_rows()
        for cell in row
    )
    matrix.add(
        surface="Monthly Excel export",
        source_table="anomalies+visits",
        source_field="yyyymm",
        service_function="event_service.export_monthly_excel",
        output_artifact=monthly_xlsx.name,
        output_location="workbook text",
        expected=True,
        actual=ok and supplier_name in monthly_text,
        evidence_path=monthly_xlsx,
    )
    matrix.add_contains(
        surface="Monthly Excel export message",
        source_table="monthly_stats_cache",
        source_field="yyyymm",
        service_function="event_service.export_monthly_excel",
        output_artifact=monthly_xlsx.name,
        output_location="return message",
        expected="已匯出至",
        actual=msg,
        evidence_path=monthly_xlsx,
    )

    events_xlsx = artifacts_dir / "event_range_report.xlsx"
    ok, msg = event_service.export_events_report(
        str(events_xlsx), args.start_date, args.end_date
    )
    range_wb = load_workbook(events_xlsx, data_only=False)
    anomaly_sheet = range_wb["異常"]
    anomaly_values = [
        [anomaly_sheet.cell(row=r, column=c).value for c in range(1, 10)]
        for r in range(1, anomaly_sheet.max_row + 1)
    ]
    matrix.add(
        surface="Range Excel anomaly split",
        source_table="anomalies",
        source_field="quality_report_required",
        service_function="event_service.export_events_report",
        output_artifact=events_xlsx.name,
        output_location="異常!I:I",
        expected=True,
        actual=ok and any(row[0] == str(standalone_no) and row[8] == "是" for row in anomaly_values),
        evidence_path=events_xlsx,
    )
    matrix.add_contains(
        surface="Range Excel export message",
        source_table="anomalies+visits",
        source_field="event_date",
        service_function="event_service.export_events_report",
        output_artifact=events_xlsx.name,
        output_location="return message",
        expected="已匯出至",
        actual=msg,
        evidence_path=events_xlsx,
    )

    standalone_detail = event_service.get_anomaly_detail(standalone_id)
    pdf_row = _row_by_value(events, "ref_no", str(standalone_no))
    full_pdf = artifacts_dir / "event_full.pdf"
    ok, msg = event_pdf_exporter.export_event_pdf(full_pdf, pdf_row, standalone_detail)
    pdf_text, pdf_status = _load_pdf_text(full_pdf)
    matrix.add(
        surface="Event PDF artifact",
        source_table="anomalies",
        source_field="problem_desc",
        service_function="event_pdf_exporter.export_event_pdf",
        output_artifact=full_pdf.name,
        output_location="PDF extracted text",
        expected=True,
        actual=ok and full_pdf.exists() and full_pdf.stat().st_size > 1000,
        evidence_path=full_pdf,
    )
    matrix.add_contains(
        surface="Event PDF text parity",
        source_table="anomalies",
        source_field="anomaly_no",
        service_function="event_pdf_exporter.export_event_pdf",
        output_artifact=full_pdf.name,
        output_location="PDF extracted text",
        expected=str(standalone_no),
        actual=pdf_text if pdf_status == "pass" else pdf_status,
        evidence_path=full_pdf,
    )

    brief_pdf = artifacts_dir / "event_brief.pdf"
    ok, msg = event_pdf_exporter.export_brief_event_pdf(brief_pdf, pdf_row, standalone_detail)
    brief_image = event_pdf_exporter.render_brief_event_to_image(pdf_row, standalone_detail)
    brief_image_path = artifacts_dir / "event_brief.png"
    image_saved = bool(brief_image and brief_image.save(str(brief_image_path)))
    matrix.add(
        surface="Brief PDF artifact",
        source_table="anomalies",
        source_field="problem_desc",
        service_function="event_pdf_exporter.export_brief_event_pdf",
        output_artifact=brief_pdf.name,
        output_location="file size",
        expected=True,
        actual=ok and brief_pdf.exists() and brief_pdf.stat().st_size > 1000,
        evidence_path=brief_pdf,
    )
    matrix.add(
        surface="Brief image artifact",
        source_table="anomalies",
        source_field="problem_desc",
        service_function="event_pdf_exporter.render_brief_event_to_image",
        output_artifact=brief_image_path.name,
        output_location="image saved",
        expected=True,
        actual=image_saved,
        evidence_path=brief_image_path,
    )

    with connection.get_connection() as conn:
        defect_no = defect_service.create_defect(
            conn,
            {
                "event_date": "2026-06-13",
                "return_slip_type": "廠內退料",
                "processing_line": "原物料",
                "work_order_no": f"NCR-WO-{stamp}",
                "internal_work_order_no": f"NCR-IWO-{stamp}",
                "transfer_slip_no": f"NCR-TS-{stamp}",
                "item_no": f"NCR-ITEM-{stamp}",
                "product_name": product_name,
                "qty": 7,
                "category": "原物料",
                "supplier_name": supplier_name,
                "outsource_supplier_name": "N/A",
                "defect_desc": f"{stamp} warehouse defect input",
                "status": "處理中",
                "disposition": "重工",
                "responsibility": "材損",
            },
        )
    with connection.get_connection() as conn:
        defects = [dict(row) for row in ncr_stats_service.get_defects_detail_by_range(conn, args.start_date, args.end_date)]
        audit_defects = [row for row in defects if row.get("defect_no") == defect_no]
        ncr_summary = ncr_stats_service.get_warehouse_nonconforming_summary(conn)
        defect_count_for_visit_note = conn.execute(
            "SELECT COUNT(*) AS c FROM defect_records WHERE defect_desc = ?",
            (f"{stamp} visit defect note input",),
        ).fetchone()["c"]
    ncr_xlsx = artifacts_dir / "ncr_report.xlsx"
    ok, msg = ncr_export_service.export_ncr_excel_report(
        str(ncr_xlsx), args.start_date, args.end_date, audit_defects
    )
    ncr_wb = load_workbook(ncr_xlsx, data_only=False)
    ncr_detail = ncr_wb["不良品明細"]
    matrix.add(
        surface="Warehouse NCR boundary",
        source_table="visit_defect_notes",
        source_field="defect_desc",
        service_function="ncr_stats_service.get_warehouse_nonconforming_summary",
        output_artifact="service payload",
        output_location="defect_records count for visit note",
        expected=0,
        actual=defect_count_for_visit_note,
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Warehouse NCR stats",
        source_table="defect_records",
        source_field="qty",
        service_function="ncr_stats_service.get_warehouse_nonconforming_summary",
        output_artifact="service payload",
        output_location="summary.total_qty >= audit qty",
        expected=True,
        actual=ncr_summary["total_qty"] >= 7,
        evidence_path=output_dir / "evidence_matrix.json",
    )
    matrix.add(
        surface="Warehouse NCR Excel detail",
        source_table="defect_records",
        source_field="defect_no",
        service_function="ncr_export_service.export_ncr_excel_report",
        output_artifact=ncr_xlsx.name,
        output_location="不良品明細!A:A",
        expected=defect_no,
        actual=ncr_detail.cell(row=2, column=1).value,
        evidence_path=ncr_xlsx,
    )

    import scripts.generate_weekly_report as weekly_report

    weekly_report.OUT_DIR = artifacts_dir
    pptx_path = weekly_report.generate_report()
    pptx_text = _pptx_text(pptx_path)
    matrix.add_contains(
        surface="Weekly PPTX report",
        source_table="anomalies",
        source_field="supplier_name",
        service_function="scripts.generate_weekly_report.generate_report",
        output_artifact=pptx_path.name,
        output_location="slides text",
        expected=supplier_name,
        actual=pptx_text,
        evidence_path=pptx_path,
    )

    summary = {
        "supplier_name": supplier_name,
        "product_code": product_code,
        "standalone_anomaly_no": str(standalone_no),
        "visit_id": visit_id,
        "confirmed_anomaly_id": confirmed.get("anomaly_id"),
        "linked_anomaly_id": linked.get("anomaly_id"),
        "defect_no": defect_no,
        "ranking_rows": ranking,
    }
    (output_dir / "audit_seed_summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return summary


def main() -> int:
    args = _parse_args()
    output_dir = _prepare_output_dir(args.output_dir)
    db_copy = _prepare_disposable_db(args.db_path.expanduser().resolve(), output_dir)
    _setup_environment(db_copy, output_dir)
    matrix = EvidenceMatrix()
    try:
        summary = _seed_and_audit(args, output_dir, matrix)
        status = "pass" if matrix.exit_code() == 0 else "not_pass"
    except Exception as exc:
        matrix.add(
            surface="Audit harness execution",
            source_table="n/a",
            source_field="n/a",
            service_function="scripts.audit_report_parity.main",
            output_artifact="evidence_matrix.json",
            output_location="exception",
            expected="no exception",
            actual=f"{type(exc).__name__}: {exc}",
            status="blocked",
            evidence_path=output_dir,
        )
        summary = {"error": f"{type(exc).__name__}: {exc}"}
        status = "blocked"

    json_path, md_path = matrix.write(output_dir)
    run_summary = {
        "status": status,
        "output_dir": str(output_dir),
        "disposable_db": str(db_copy),
        "evidence_json": str(json_path),
        "evidence_markdown": str(md_path),
        "seed_summary": summary,
        "rows": len(matrix.rows),
        "not_pass": [asdict(row) for row in matrix.rows if row.status != "pass"],
    }
    summary_path = output_dir / "audit_summary.json"
    summary_path.write_text(json.dumps(run_summary, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(run_summary, ensure_ascii=False, indent=2))
    return 0 if status == "pass" else 1


if __name__ == "__main__":
    raise SystemExit(main())
