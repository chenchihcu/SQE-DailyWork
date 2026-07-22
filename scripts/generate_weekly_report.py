"""
SQE Weekly Report Generator
產出: Outputs/SQE_<YYYY>W<WW>.pptx
資料來源: data/sqe_v2.db
"""

import sqlite3
import sys
from datetime import date, datetime, timedelta
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ── 路徑設定 ──────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent.parent
SRC_DIR = BASE_DIR / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from database.connection import DB_PATH  # noqa: E402

OUT_DIR  = BASE_DIR / "Outputs"

# ── Mitcorp 品牌色系 ──────────────────────────────────────
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
C_NAVY       = RGBColor(0x0D, 0x20, 0x3E)
C_STEEL      = RGBColor(0x1E, 0x3C, 0x56)
C_FOOTER     = RGBColor(0x2D, 0x4A, 0x6A)
C_GREEN_MC   = RGBColor(0x76, 0xB8, 0x2A)
C_TEAL_MC    = RGBColor(0x00, 0x9A, 0x96)
C_GRAY_TOP   = RGBColor(0xE8, 0xEA, 0xED)
C_HEADER_BG  = RGBColor(0x2D, 0x4A, 0x6A)
C_ALT_ROW    = RGBColor(0xF0, 0xF4, 0xF8)
C_ACCENT     = RGBColor(0xC0, 0x39, 0x2B)
C_BORDER     = RGBColor(0xCC, 0xD3, 0xDE)
C_NEW_ITEM   = RGBColor(0xD4, 0x7B, 0x00)   # 琥珀橘，本週新增高亮

FONT_NAME    = "微軟正黑體"
FONT_NAME_EN = "Calibri"


SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)

# ── 列表頁分頁參數 ────────────────────────────────────────
ROWS_PER_PAGE     = 14    # 每頁最多列數,預留底部空間
TABLE_TOP_CM      = 4.0
HEADER_H_CM       = 0.75
ROW_H_CM          = 0.65
BOTTOM_RESERVE_CM = 1.2   # 底部與 footer 留白

# ── 資料層 ────────────────────────────────────────────────

def get_conn():
    conn = sqlite3.connect(f"file:{DB_PATH.as_posix()}?mode=ro", uri=True)
    conn.row_factory = sqlite3.Row
    return conn

def fetch_open_anomalies(conn):
    cur = conn.cursor()
    cur.execute("""
        SELECT a.anomaly_no, a.anomaly_date, a.category, a.product_name,
               a.product_stage, a.problem_desc, a.pending_items, a.improvement_desc,
               a.responsible_person, a.due_date, a.root_cause_category, a.status,
               a.visit_id,
               s.supplier_name
        FROM anomalies a
        LEFT JOIN suppliers s ON a.supplier_id = s.id
        WHERE a.status != '已結案'
        ORDER BY s.supplier_name, a.anomaly_date DESC
    """)
    return [dict(r) for r in cur.fetchall()]

def fetch_visits(conn):
    cur = conn.cursor()
    cur.execute("""
        SELECT v.id, v.visit_date, v.product_name, v.product_stage,
               v.work_order_no, v.summary, v.status, s.supplier_name
        FROM visits v
        LEFT JOIN suppliers s ON v.supplier_id = s.id
        ORDER BY v.visit_date DESC
    """)
    return [dict(r) for r in cur.fetchall()]

def fetch_visit_anomalies(conn):
    cur = conn.cursor()
    cur.execute("""
        SELECT a.anomaly_no, a.status, a.problem_desc, a.pending_items,
               a.improvement_desc, a.responsible_person, a.due_date, a.visit_id,
               s.supplier_name
        FROM anomalies a
        LEFT JOIN suppliers s ON a.supplier_id = s.id
        WHERE a.visit_id IS NOT NULL AND a.visit_id != ''
          AND a.status != '已結案'
    """)
    result = {}
    for row in cur.fetchall():
        d = dict(row)
        vid = d["visit_id"]
        result.setdefault(vid, []).append(d)
    return result

# ── 通用 Helper ───────────────────────────────────────────

def trunc(text: str, max_len: int) -> str:
    if not text:
        return ""
    text = text.replace("\n", " ").strip()
    return text if len(text) <= max_len else text[:max_len] + "…"

_TODAY = date.today()

def week_label() -> str:
    year, week, _ = _TODAY.isocalendar()
    return f"{year}W{week:02d}"

_CN_NUMS = "一二三四五六七八九十"
def _cn(n: int) -> str:
    """1→'一', 2→'二', ...；超過 10 退化為阿拉伯數字。"""
    return _CN_NUMS[n-1] if 1 <= n <= len(_CN_NUMS) else str(n)

def set_cell_bg(cell, rgb: RGBColor):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")

def set_cell_border(cell, rgb: RGBColor, width_pt: float = 0.5):
    def make_ln(parent_tag):
        ln = etree.Element(qn(parent_tag))
        ln.set("w", str(int(width_pt * 12700)))
        sf = etree.SubElement(ln, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
        return ln
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for old in tcPr.findall(qn(side)):
            tcPr.remove(old)
        tcPr.append(make_ln(side))

def add_cell_text(cell, text: str, bold: bool = False, color: RGBColor = C_BLACK,
                  size_pt: float = 9, align=PP_ALIGN.LEFT):
    tf  = cell.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    run.font.bold  = bold
    run.font.size  = Pt(size_pt)
    run.font.color.rgb = color
    run.font.name  = FONT_NAME

def style_header_row(table, col_count: int):
    for ci in range(col_count):
        cell = table.cell(0, ci)
        set_cell_bg(cell, C_HEADER_BG)
        set_cell_border(cell, C_WHITE, 0.5)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = C_WHITE
                run.font.bold      = True
                run.font.size      = Pt(9)
                run.font.name      = FONT_NAME

def style_data_rows(table, row_count: int, col_count: int):
    for ri in range(1, row_count):
        bg = C_ALT_ROW if ri % 2 == 0 else C_WHITE
        for ci in range(col_count):
            cell = table.cell(ri, ci)
            set_cell_bg(cell, bg)
            set_cell_border(cell, C_BORDER, 0.3)

def calc_col_widths(headers, data, total_width_cm):
    """根據資料長度自動分配欄寬"""
    def char_w(text):
        # 中文字元計算為 2 寬度
        text = str(text or "")
        return sum(2 if ord(c) > 127 else 1 for c in text)
    
    max_lens = [max(10, char_w(h)) for h in headers]
    for row in data:
        for i, val in enumerate(row):
            max_lens[i] = max(max_lens[i], char_w(val))
    
    total_len = sum(max_lens)
    if total_len == 0:
        return [Cm(total_width_cm / len(headers))] * len(headers)
    
    return [Cm(total_width_cm * (l / total_len)) for l in max_lens]

# ── 投影片建構 ────────────────────────────────────────────

def add_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def _bezier_sample(p0, p1, p2, p3, n=20):
    pts = []
    for i in range(n + 1):
        t = i / n
        x = (1-t)**3*p0[0] + 3*(1-t)**2*t*p1[0] + 3*(1-t)*t**2*p2[0] + t**3*p3[0]
        y = (1-t)**3*p0[1] + 3*(1-t)**2*t*p1[1] + 3*(1-t)*t**2*p2[1] + t**3*p3[1]
        pts.append((x, y))
    return pts

def _make_polygon(slide, pts_cm, color: RGBColor):
    if len(pts_cm) < 2:
        return
    builder = slide.shapes.build_freeform(Cm(pts_cm[0][0]), Cm(pts_cm[0][1]))
    builder.add_line_segments([(Cm(x), Cm(y)) for x, y in pts_cm[1:]])
    shape = builder.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _add_logo(slide, left_cm, top_cm, height_cm):
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(5), Cm(height_cm))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Mitcorp"
    r.font.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = C_WHITE
    r.font.name = FONT_NAME_EN

def add_content_bg(slide, title: str, subtitle: str = ""):
    W_cm = 33.867
    H_cm = 19.05
    FOOTER_H = 1.85
    TOP_H = 3.8

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_WHITE
    bg.line.fill.background()

    top = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Cm(TOP_H))
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(0xEA, 0xEC, 0xEF)
    top.line.fill.background()

    sep = slide.shapes.add_shape(1, 0, Cm(TOP_H), SLIDE_W, Cm(0.07))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xBB, 0xC2, 0xCC)
    sep.line.fill.background()

    footer = slide.shapes.add_shape(1, 0, Cm(H_cm - FOOTER_H), SLIDE_W, Cm(FOOTER_H))
    footer.fill.solid()
    footer.fill.fore_color.rgb = C_FOOTER
    footer.line.fill.background()

    _add_logo(slide, 0.5, H_cm - FOOTER_H + 0.3, FOOTER_H * 0.65)

    tb_ir = slide.shapes.add_textbox(Cm(W_cm - 8.0), Cm(H_cm - FOOTER_H + 0.55), Cm(7.6), Cm(0.9))
    tf = tb_ir.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "INSPECTION REDEFINED"
    run.font.size  = Pt(10)
    run.font.bold  = True
    run.font.color.rgb = C_WHITE
    run.font.name  = FONT_NAME_EN

    tb_t = slide.shapes.add_textbox(Cm(0.8), Cm(0.6), Cm(30), Cm(2.0))
    tf_t = tb_t.text_frame
    p_t  = tf_t.paragraphs[0]
    r_t  = p_t.add_run()
    r_t.text = title
    r_t.font.bold  = True
    r_t.font.size  = Pt(20)
    r_t.font.color.rgb = C_NAVY
    r_t.font.name  = FONT_NAME

    if subtitle:
        p_s = tf_t.add_paragraph()
        r_s = p_s.add_run()
        r_s.text = subtitle
        r_s.font.size  = Pt(10)
        r_s.font.color.rgb = C_FOOTER
        r_s.font.name  = FONT_NAME

def add_bullet_box(slide, lines: list[str], left, top, width, height,
                   color: RGBColor = C_BLACK, size_pt: float = 9):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        run.font.size  = Pt(size_pt)
        run.font.color.rgb = color
        run.font.name  = FONT_NAME

# ── Slide 1：封面 ────────────────────────

def build_cover(prs: Presentation):
    slide = add_slide(prs)
    W = 33.867
    H = 19.05

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_NAVY
    bg.line.fill.background()

    GREEN_T = 1.2
    TEAL_T  = 1.0
    P0 = (0,   11.4)
    P1 = (8,   11.2)
    P2 = (25,  15.5)
    P3 = (W,   16.2)
    center = _bezier_sample(P0, P1, P2, P3, n=28)

    steel_pts = [(x, y + TEAL_T) for x, y in center] + [(W, H), (0, H)]
    _make_polygon(slide, steel_pts, C_STEEL)

    teal_pts = list(center) + list(reversed([(x, y + TEAL_T) for x, y in center]))
    _make_polygon(slide, teal_pts, C_TEAL_MC)

    green_pts = [(x, y - GREEN_T) for x, y in center] + list(reversed(center))
    _make_polygon(slide, green_pts, C_GREEN_MC)

    _add_logo(slide, 24.5, 17.1, 1.8)

    tb_ir = slide.shapes.add_textbox(Cm(24.0), Cm(18.0), Cm(9.5), Cm(0.8))
    tf_ir = tb_ir.text_frame
    p_ir  = tf_ir.paragraphs[0]
    r_ir  = p_ir.add_run()
    r_ir.text = "INSPECTION REDEFINED"
    r_ir.font.size  = Pt(9)
    r_ir.font.bold  = True
    r_ir.font.color.rgb = C_WHITE
    r_ir.font.name  = FONT_NAME_EN

    tb = slide.shapes.add_textbox(Cm(2.5), Cm(2.5), Cm(22), Cm(2.5))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = "SQE 週會報告"
    r.font.bold  = True
    r.font.size  = Pt(34)
    r.font.color.rgb = C_WHITE
    r.font.name  = FONT_NAME

    tb2 = slide.shapes.add_textbox(Cm(2.5), Cm(5.3), Cm(22), Cm(1.2))
    tf2 = tb2.text_frame
    p2  = tf2.paragraphs[0]
    r2  = p2.add_run()
    r2.text = f"週次：{week_label()}   ／   報告日期：{_TODAY.strftime('%Y/%m/%d')}"
    r2.font.size  = Pt(13)
    r2.font.color.rgb = C_TEAL_MC
    r2.font.name  = FONT_NAME

# ── 矩陣表產生器 ──────────────────────────────────────────────

def _paginate(data: list, page_size: int = ROWS_PER_PAGE) -> list[list]:
    """將資料切成等長頁面;空資料回傳 [[]] 以確保仍會建立一張投影片。"""
    if not data:
        return [[]]
    return [data[i:i + page_size] for i in range(0, len(data), page_size)]


def add_matrix_table(slide, title, cols_def, data_dicts,
                     top_cm: float = TABLE_TOP_CM,
                     new_within_days: int = 0):
    """
    單頁矩陣表生成。呼叫端負責資料分頁。
    cols_def: [(標題, 欄位key, 截斷長度)]
    data_dicts: 已切好的「當頁」資料清單
    new_within_days: > 0 時，anomaly_date 在此天數內的列以琥珀色高亮
    """
    if not data_dicts:
        add_bullet_box(slide, ["目前無資料"], Cm(1), Cm(top_cm), Cm(30), Cm(2), color=C_NAVY, size_pt=12)
        return

    HEADER_H = Cm(HEADER_H_CM)
    ROW_H    = Cm(ROW_H_CM)
    LEFT     = Cm(0.3)
    total_w  = 33.2  # 滿版寬度
    today_str  = _TODAY.isoformat()
    cutoff_str = (_TODAY - timedelta(days=new_within_days)).isoformat() \
                 if new_within_days > 0 else None

    headers = [c[0] for c in cols_def]
    keys = [c[1] for c in cols_def]
    trunc_lens = [c[2] for c in cols_def]

    table_data = []
    for row in data_dicts:
        r_data = []
        for key, tlen in zip(keys, trunc_lens, strict=True):
            val = str(row.get(key) or "")
            if tlen > 0:
                val = trunc(val, tlen)
            r_data.append(val)
        table_data.append(r_data)

    col_widths = calc_col_widths(headers, table_data, total_w)

    row_count = len(table_data)
    table = slide.shapes.add_table(
        row_count + 1,
        len(cols_def),
        LEFT, Cm(top_cm), sum(col_widths),
        HEADER_H + ROW_H * row_count
    ).table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ci, label in enumerate(headers):
        add_cell_text(table.cell(0, ci), label, bold=True, color=C_WHITE, size_pt=9, align=PP_ALIGN.CENTER)

    style_header_row(table, len(cols_def))

    for ri, (r_data, raw_row) in enumerate(zip(table_data, data_dicts, strict=True), start=1):
        is_overdue  = False
        is_new_item = False

        if "due_date" in raw_row:
            due = raw_row.get("due_date") or ""
            if due and due != "未設定" and due < today_str:
                is_overdue = True

        if cutoff_str and "anomaly_date" in raw_row:
            adate = raw_row.get("anomaly_date") or ""
            if adate and adate >= cutoff_str:
                is_new_item = True

        for ci, val in enumerate(r_data):
            if is_overdue:
                color = C_ACCENT       # 優先：逾期紅
            elif is_new_item:
                color = C_NEW_ITEM     # 次優：本週新增琥珀
            else:
                color = C_BLACK
            add_cell_text(table.cell(ri, ci), val, size_pt=8.5, color=color)

    style_data_rows(table, row_count + 1, len(cols_def))

# ── Slide 2：KPI 統計列 ───────────────────────────────────

def add_kpi_row(slide, anomalies: list[dict]) -> None:
    """在 slide 頂部繪製 4 個 KPI 統計方塊（僅第一頁呼叫）。"""
    today_str  = _TODAY.isoformat()
    cutoff_str = (_TODAY - timedelta(days=7)).isoformat()

    total     = len(anomalies)
    overdue   = sum(1 for a in anomalies
                    if (a.get("due_date") or "") not in ("", "未設定")
                    and a["due_date"] < today_str)
    suppliers = len({a["supplier_name"] for a in anomalies if a.get("supplier_name")})
    new_7d    = sum(1 for a in anomalies
                    if (a.get("anomaly_date") or "") >= cutoff_str)

    kpis = [
        ("未結案總數", total,     C_NAVY),
        ("逾期筆數",   overdue,   C_ACCENT   if overdue > 0 else C_NAVY),
        ("涉及廠商數", suppliers, C_NAVY),
        ("本週新增",   new_7d,    C_NEW_ITEM if new_7d  > 0 else C_NAVY),
    ]

    BOX_TOP = Cm(2.65)
    BOX_H   = Cm(2.5)
    BOX_W   = Cm(8.1)
    GAP     = Cm(0.25)
    LEFT0   = Cm(0.3)

    for i, (label, value, num_color) in enumerate(kpis):
        left = LEFT0 + i * (BOX_W + GAP)

        box = slide.shapes.add_shape(1, left, BOX_TOP, BOX_W, BOX_H)
        box.fill.solid()
        box.fill.fore_color.rgb = C_GRAY_TOP
        box.line.fill.background()

        tb_n = slide.shapes.add_textbox(
            left + Cm(0.2), BOX_TOP + Cm(0.3), BOX_W - Cm(0.4), Cm(1.5))
        p = tb_n.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(value)
        run.font.size      = Pt(22)
        run.font.bold      = True
        run.font.color.rgb = num_color
        run.font.name      = FONT_NAME_EN

        tb_l = slide.shapes.add_textbox(
            left + Cm(0.2), BOX_TOP + Cm(1.8), BOX_W - Cm(0.4), Cm(0.7))
        p2 = tb_l.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text           = label
        run2.font.size      = Pt(9)
        run2.font.color.rgb = C_FOOTER
        run2.font.name      = FONT_NAME

# ── Slide 2：待處理異常總覽（分類三頁） ───────────────────

def categorize_anomalies(anomalies: list[dict]) -> tuple[list, list, list]:
    """
    將未結案異常依互斥規則分成三類（逾期優先）：
      - new_items:     7 天內 anomaly_date，且未逾期
      - overdue_items: due_date < today（且 due_date 有設定）
      - others:        其餘
    輸入排序保留（呼叫端負責 supplier_name 排序）。
    """
    today_str  = _TODAY.isoformat()
    cutoff_str = (_TODAY - timedelta(days=7)).isoformat()

    new_items, overdue_items, others = [], [], []
    for a in anomalies:
        due = a.get("due_date") or ""
        is_overdue = bool(due and due != "未設定" and due < today_str)

        adate = a.get("anomaly_date") or ""
        is_new = bool(adate and adate >= cutoff_str)

        if is_overdue:
            overdue_items.append(a)
        elif is_new:
            new_items.append(a)
        else:
            others.append(a)

    return new_items, overdue_items, others


COLS_FULL = [   # 含截止日：逾期事件 / 其他事件
    ("日期",     "anomaly_date",       0),
    ("來源",     "_source",            0),
    ("廠商",     "supplier_name",      0),
    ("產品",     "product_name",       30),
    ("問題描述", "problem_desc",       60),
    ("負責人",   "responsible_person", 0),
    ("截止日",   "due_date",           0),
]

COLS_NO_DUE = [  # 移除截止日：本週新增
    ("日期",     "anomaly_date",       0),
    ("來源",     "_source",            0),
    ("廠商",     "supplier_name",      0),
    ("產品",     "product_name",       30),
    ("問題描述", "problem_desc",       60),
    ("負責人",   "responsible_person", 0),
]


def _tag_source(anomalies: list[dict]) -> None:
    """In-place 標註每筆異常的來源：● 訪廠 / ○ 獨立。"""
    for a in anomalies:
        vid = a.get("visit_id")
        a["_source"] = "● 訪廠" if (vid is not None and str(vid).strip()) else "○ 獨立"


def _emit_category_pages(prs, all_anomalies, items, section_label, cols,
                         with_kpi: bool, new_within_days: int,
                         chapter_no: int):
    """為單一分類產生 1~N 頁。空集合輸出單張無資料頁。"""
    FIRST_KPI_ROWS = 10
    FIRST_KPI_TOP  = 5.5

    if items:
        if with_kpi:
            first_chunk = items[:FIRST_KPI_ROWS]
            remaining   = items[FIRST_KPI_ROWS:]
            rest        = _paginate(remaining) if remaining else []
            pages       = [first_chunk] + rest
        else:
            pages = _paginate(items)
    else:
        pages = [[]]

    total = len(pages)
    for idx, chunk in enumerate(pages, 1):
        slide    = add_slide(prs)
        page_tag = f" ({idx}/{total})" if total > 1 else ""
        add_content_bg(slide,
                       f"{_cn(chapter_no)}. 待處理異常 — {section_label}{page_tag}",
                       f"{section_label} 共 {len(items)} 筆")

        if with_kpi and idx == 1:
            add_kpi_row(slide, all_anomalies)
            top_cm = FIRST_KPI_TOP
        else:
            top_cm = TABLE_TOP_CM

        add_matrix_table(slide, section_label, cols, chunk,
                         top_cm=top_cm, new_within_days=new_within_days)


def build_anomaly_slide(prs: Presentation, anomalies: list[dict],
                        chapter_start: int = 1) -> int:
    """產生待處理異常的 3 個分類章節，回傳下一個可用章節號。"""
    _tag_source(anomalies)
    new_items, overdue_items, others = categorize_anomalies(anomalies)

    # #1 本週新增（KPI + 無截止日欄；items 全為非逾期，套琥珀高亮）
    _emit_category_pages(prs, anomalies, new_items,
                         section_label="本週新增",
                         cols=COLS_NO_DUE,
                         with_kpi=True,
                         new_within_days=7,
                         chapter_no=chapter_start)

    # #2 逾期事件（含截止日；items 全為逾期，套紅色高亮）
    _emit_category_pages(prs, anomalies, overdue_items,
                         section_label="逾期事件",
                         cols=COLS_FULL,
                         with_kpi=False,
                         new_within_days=0,
                         chapter_no=chapter_start + 1)

    # #3 其他事件（含截止日；無高亮）
    _emit_category_pages(prs, anomalies, others,
                         section_label="其他事件",
                         cols=COLS_FULL,
                         with_kpi=False,
                         new_within_days=0,
                         chapter_no=chapter_start + 2)

    return chapter_start + 3

# ── Slide 3：訪廠紀錄 (無異常) ─────────────────────────

def build_visit_normal_slide(prs: Presentation, visits: list[dict], v_anoms: dict,
                             chapter_no: int = 1) -> int:
    # 僅顯示本週（近 7 天）且未發生異常的訪廠；歷史資料屬於過去式，無討論價值
    cutoff_str = (_TODAY - timedelta(days=7)).isoformat()
    normal_visits = [v for v in visits
                     if not v_anoms.get(v.get("id"))
                     and (v.get("visit_date") or "") >= cutoff_str]
    cols = [
        ("訪廠日期", "visit_date", 0),
        ("廠商", "supplier_name", 0),
        ("產品", "product_name", 30),
        ("產品階段", "product_stage", 0),
        ("工單號", "work_order_no", 0),
        ("摘要", "summary", 60),
    ]
    pages = _paginate(normal_visits)
    total_pages = len(pages)
    for idx, chunk in enumerate(pages, 1):
        slide = add_slide(prs)
        page_tag = f" ({idx}/{total_pages})" if total_pages > 1 else ""
        add_content_bg(slide, f"{_cn(chapter_no)}. 訪廠紀錄 (未發生異常){page_tag}",
                       f"本週共 {len(normal_visits)} 筆 (近 7 天)")
        add_matrix_table(slide, "訪廠紀錄", cols, chunk)
    return chapter_no + 1

# ── Slide 4：訪廠發現異常 ─────────────────────────

# ── 主流程 ────────────────────────────────────────────────

def generate_report() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    label    = week_label()
    out_path = OUT_DIR / f"SQE_{label}.pptx"

    conn      = get_conn()
    anomalies = fetch_open_anomalies(conn)
    visits    = fetch_visits(conn)
    v_anoms   = fetch_visit_anomalies(conn)
    conn.close()

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs)
    ch = 1
    # 待處理異常含「來源」欄（● 訪廠 / ○ 獨立），不再單獨產生來源視角頁
    ch = build_anomaly_slide(prs, anomalies, ch)
    ch = build_visit_normal_slide(prs, visits, v_anoms, ch)

    try:
        prs.save(str(out_path))
    except PermissionError:
        ts = datetime.now().strftime("%H%M%S")
        alt_path = out_path.with_stem(f"{out_path.stem}_{ts}")
        try:
            prs.save(str(alt_path))
            out_path = alt_path
        except PermissionError as exc:
            raise PermissionError(
                f"檔案已在其他程式中開啟，請關閉後重試。\n({out_path.name})"
            ) from exc
    return out_path

def main():
    path = generate_report()
    sys.stdout.reconfigure(encoding="utf-8")
    print(f"OK: {path}")

if __name__ == "__main__":
    main()
